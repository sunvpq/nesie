"""Generate Nesie pitch deck as PowerPoint."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BLACK = RGBColor(0x0A, 0x0A, 0x0A)
CARD = RGBColor(0x18, 0x18, 0x18)
GREEN = RGBColor(0xAA, 0xFF, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x99, 0x99, 0x99)
RED = RGBColor(0xFF, 0x44, 0x44)
YELLOW = RGBColor(0xFF, 0xCC, 0x00)
DARK_GREEN = RGBColor(0x1A, 0x2A, 0x00)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# --- Helpers ---

def set_bg(slide, color=BLACK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height,
             font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_para(text_frame, text, font_size=16, color=WHITE, bold=False,
             alignment=PP_ALIGN.LEFT, space_before=Pt(6), font_name="Calibri"):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    return p


def add_rounded_rect(slide, left, top, width, height, fill_color=CARD, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_slide_number(slide, num):
    add_text(slide, f"{num} / 10",
             W - Inches(1.2), H - Inches(0.6), Inches(1), Inches(0.4),
             font_size=11, color=GRAY, alignment=PP_ALIGN.RIGHT)


def card_with_text(slide, left, top, w, h, lines, fill=CARD, border=None):
    """Add a rounded rect card with multiple text lines.
    lines: list of (text, font_size, color, bold)
    """
    add_rounded_rect(slide, left, top, w, h, fill, border)
    pad = Inches(0.25)
    txBox = slide.shapes.add_textbox(left + pad, top + pad, w - 2*pad, h - 2*pad)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, fs, clr, bld) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(6)
        p.text = text
        p.font.size = Pt(fs)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = "Calibri"
    return txBox


# ============================================================
# SLIDE 1 — COVER
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(s)

# Green circle behind logo
cx, cy = W // 2, H // 2
r = Inches(2.2)
circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                            cx - r, cy - r - Inches(0.3), r*2, r*2)
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x00)
circle.line.fill.background()

add_text(s, "Nesie",
         Inches(0), Inches(2.0), W, Inches(1.5),
         font_size=96, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s, "Кредитный рейтинг под контролем",
         Inches(0), Inches(3.7), W, Inches(0.7),
         font_size=28, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text(s, "Финтех стартап  ·  Атырау  ·  2026",
         Inches(0), Inches(4.5), W, Inches(0.5),
         font_size=16, color=GREEN, alignment=PP_ALIGN.CENTER)
add_slide_number(s, 1)

# ============================================================
# SLIDE 2 — STATISTICS
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, "Казахстан в долгах",
         Inches(0), Inches(0.5), W, Inches(0.8),
         font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

stats = [
    ("10.2 млн", "активных заёмщиков"),
    ("2.3", "кредита на одного человека"),
    ("14.7 трлн ₸", "объём кредитования"),
    ("34%", "казахстанцев — долговая\nнагрузка выше 50% дохода"),
]
card_w = Inches(2.8)
card_h = Inches(3.2)
gap = Inches(0.3)
total_w = 4 * card_w + 3 * gap
start_x = (W - total_w) // 2
y = Inches(2.0)

for i, (num, label) in enumerate(stats):
    x = start_x + i * (card_w + gap)
    card_with_text(s, x, y, card_w, card_h, [
        (num, 40, GREEN, True),
        (label, 16, GRAY, False),
    ])

add_slide_number(s, 2)

# ============================================================
# SLIDE 3 — HUMAN STORY
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, "История Адила, 26 лет, Атырау",
         Inches(0), Inches(0.4), W, Inches(0.8),
         font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

timeline = [
    ("Январь", "Взял кредит на телефон", "Скор: 680 ✓", GREEN),
    ("Март", "Помог сестре, ещё кредит", "Скор: 650 (не знает)", YELLOW),
    ("Июнь", "Кредит на холодильник", "Скор: 590 (не знает)", YELLOW),
    ("Август", "Просрочил на 12 дней", "Скор: 530 (не знает)", RED),
    ("Ноябрь", "Подал на автокредит", "ОТКАЗ ✗", RED),
    ("Ноябрь", "МФО под 98% годовых", "Долговая ловушка", RED),
]

start_y = Inches(1.5)
row_h = Inches(0.85)
left_col = Inches(2.5)
mid_col = Inches(4.2)
right_col = Inches(8.5)

for i, (month, desc, score, clr) in enumerate(timeline):
    y = start_y + i * row_h
    # Dot
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             Inches(3.8), y + Inches(0.08), Inches(0.22), Inches(0.22))
    dot.fill.solid()
    dot.fill.fore_color.rgb = clr
    dot.line.fill.background()
    # Line
    if i < len(timeline) - 1:
        line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(3.88), y + Inches(0.32), Pt(3), row_h - Inches(0.1))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
        line.line.fill.background()
    # Month
    add_text(s, month, left_col, y, Inches(1.5), Inches(0.4),
             font_size=16, color=WHITE, bold=True)
    # Description
    add_text(s, desc, mid_col, y, Inches(4), Inches(0.4),
             font_size=16, color=GRAY)
    # Score
    add_text(s, score, right_col, y, Inches(3.5), Inches(0.4),
             font_size=16, color=clr, bold=True)

add_slide_number(s, 3)

# ============================================================
# SLIDE 4 — ROOT CAUSE
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, "Люди узнают о проблеме когда уже поздно",
         Inches(0), Inches(0.5), W, Inches(0.8),
         font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

problems = [
    ("🔒  Скор существует — но скрыт",
     "ПКБ берёт 400₸ и требует ручного запроса через неудобный сайт"),
    ("🏦  Кредиты в разных банках — нет единой картины",
     "Человек не видит общую нагрузку и не понимает совокупный риск"),
    ("⚠️  Никто не предупреждает до оформления",
     "Решение о новом кредите принимается вслепую, без данных о последствиях"),
]

card_w2 = Inches(10)
card_h2 = Inches(1.4)
start_x2 = (W - card_w2) // 2
start_y2 = Inches(2.0)

for i, (title, desc) in enumerate(problems):
    y = start_y2 + i * (card_h2 + Inches(0.25))
    card_with_text(s, start_x2, y, card_w2, card_h2, [
        (title, 22, WHITE, True),
        (desc, 16, GRAY, False),
    ])

add_slide_number(s, 4)

# ============================================================
# SLIDE 5 — SOLUTION
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

add_text(s, "Nesie",
         Inches(0), Inches(0.8), W, Inches(1.2),
         font_size=80, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s, "Первое приложение, которое показывает казахстанцам\nих реальное финансовое здоровье",
         Inches(1.5), Inches(2.3), W - Inches(3), Inches(1.0),
         font_size=22, color=GRAY, alignment=PP_ALIGN.CENTER)

pillars = [
    ("Живой скор из ПКБ", "Актуальный кредитный рейтинг\nв одно касание"),
    ("Все кредиты в одном месте", "Единая картина долговой\nнагрузки со всех банков"),
    ("Симулятор последствий", "Увидь как изменится скор\nдо оформления кредита"),
]
pw = Inches(3.5)
ph = Inches(2.8)
pgap = Inches(0.4)
ptotal = 3 * pw + 2 * pgap
pstart = (W - ptotal) // 2

for i, (title, desc) in enumerate(pillars):
    x = pstart + i * (pw + pgap)
    card_with_text(s, x, Inches(3.8), pw, ph, [
        (title, 20, GREEN, True),
        ("", 8, GRAY, False),
        (desc, 15, GRAY, False),
    ])

add_slide_number(s, 5)

# ============================================================
# SLIDE 6 — PRODUCT DEMO
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, "Как это выглядит",
         Inches(0), Inches(0.4), W, Inches(0.7),
         font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

phone_w = Inches(3.0)
phone_h = Inches(5.5)
phone_gap = Inches(0.5)
phone_total = 3 * phone_w + 2 * phone_gap
phone_start = (W - phone_total) // 2
phone_y = Inches(1.4)

phone_data = [
    ("Скор", [
        ("714", 48, GREEN, True),
        ("", 4, GRAY, False),
        ("Хороший рейтинг", 16, GREEN, False),
        ("", 12, GRAY, False),
        ("Платежи вовремя      96%", 13, GRAY, False),
        ("Нагрузка                  42%", 13, YELLOW, False),
        ("Активных кредитов     3", 13, GRAY, False),
    ]),
    ("Мои кредиты", [
        ("Kaspi — телефон", 14, WHITE, False),
        ("142 000 ₸", 14, GREEN, True),
        ("", 4, GRAY, False),
        ("Halyk — наличные", 14, WHITE, False),
        ("580 000 ₸", 14, GREEN, True),
        ("", 4, GRAY, False),
        ("Forte — холодильник", 14, WHITE, False),
        ("215 000 ₸", 14, GREEN, True),
        ("", 8, GRAY, False),
        ("Итого: 937 000 ₸", 18, GREEN, True),
    ]),
    ("Симулятор", [
        ("Новый кредит: 500 000 ₸", 13, GRAY, False),
        ("", 8, GRAY, False),
        ("СЕЙЧАС", 11, GRAY, False),
        ("714", 40, GREEN, True),
        ("↓", 24, GRAY, True),
        ("ПОСЛЕ", 11, GRAY, False),
        ("658", 40, YELLOW, True),
        ("", 6, GRAY, False),
        ("⚠ Скор упадёт на 56 пунктов", 12, RED, True),
    ]),
]

for i, (label, lines) in enumerate(phone_data):
    x = phone_start + i * (phone_w + phone_gap)
    # Phone frame
    add_rounded_rect(s, x, phone_y, phone_w, phone_h,
                     RGBColor(0x0F, 0x0F, 0x0F), RGBColor(0x44, 0x44, 0x44))
    # Notch
    notch_w = Inches(1.0)
    add_rounded_rect(s, x + (phone_w - notch_w) // 2, phone_y,
                     notch_w, Inches(0.22),
                     RGBColor(0x0A, 0x0A, 0x0A), RGBColor(0x44, 0x44, 0x44))
    # Screen content
    pad = Inches(0.3)
    txBox = s.shapes.add_textbox(x + pad, phone_y + Inches(0.45),
                                 phone_w - 2 * pad, phone_h - Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, (text, fs, clr, bld) in enumerate(lines):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
    # Label under phone
    add_text(s, label, x, phone_y + phone_h + Inches(0.1), phone_w, Inches(0.4),
             font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(s, 6)

# ============================================================
# SLIDE 7 — COMPETITIVE ADVANTAGE
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, "Чего нет ни у кого",
         Inches(0), Inches(0.5), W, Inches(0.8),
         font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Table
table_data = [
    ["", "Kaspi", "ПКБ сайт", "Nesie"],
    ["Живой скор", "✗", "платно", "✓"],
    ["Все банки", "✗", "✓", "✓"],
    ["Симулятор", "✗", "✗", "✓"],
    ["AI советник", "✗", "✗", "✓"],
    ["Бесплатно", "✗", "✗", "✓"],
]
rows = len(table_data)
cols = 4
tw = Inches(9)
th = Inches(4.2)
tx = (W - tw) // 2
ty = Inches(1.8)

table_shape = s.shapes.add_table(rows, cols, tx, ty, tw, th)
table = table_shape.table

# Column widths
table.columns[0].width = Inches(2.8)
table.columns[1].width = Inches(2.0)
table.columns[2].width = Inches(2.0)
table.columns[3].width = Inches(2.2)

for r in range(rows):
    for c in range(cols):
        cell = table.cell(r, c)
        cell.text = table_data[r][c]
        # Styling
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.size = Pt(18 if r == 0 else 20)
                run.font.name = "Calibri"
                run.font.bold = (r == 0)
                # Color logic
                val = table_data[r][c]
                if c == 3 and r > 0:  # Nesie column
                    run.font.color.rgb = GREEN
                elif val == "✗":
                    run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
                elif val == "платно":
                    run.font.color.rgb = YELLOW
                elif val == "✓" and c != 3:
                    run.font.color.rgb = GREEN
                elif r == 0:
                    if c == 3:
                        run.font.color.rgb = GREEN
                    else:
                        run.font.color.rgb = GRAY
                else:
                    run.font.color.rgb = GRAY

        # Cell background
        cell_fill = cell.fill
        cell_fill.solid()
        if c == 3:
            cell_fill.fore_color.rgb = RGBColor(0x12, 0x1A, 0x05)
        elif r == 0:
            cell_fill.fore_color.rgb = RGBColor(0x15, 0x15, 0x15)
        else:
            cell_fill.fore_color.rgb = RGBColor(0x0F, 0x0F, 0x0F)

        cell.margin_top = Pt(8)
        cell.margin_bottom = Pt(8)

add_slide_number(s, 7)

# ============================================================
# SLIDE 8 — BUSINESS MODEL
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, "Как Nesie зарабатывает",
         Inches(0), Inches(0.5), W, Inches(0.8),
         font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

rev_cards = [
    ("💎 Подписка Pro", "990 ₸/месяц", "Симулятор, уведомления,\nAI советник"),
    ("🏦 Лидогенерация", "2 000–8 000 ₸ за лид", "Банки платят за\npre-qualified заёмщиков"),
    ("📈 Кредитный продукт", "Year 2", "Лучшая скоринговая\nмодель в Казахстане"),
]

rw = Inches(3.6)
rh = Inches(3.0)
rgap = Inches(0.35)
rtotal = 3 * rw + 2 * rgap
rstart = (W - rtotal) // 2

for i, (title, price, desc) in enumerate(rev_cards):
    x = rstart + i * (rw + rgap)
    card_with_text(s, x, Inches(1.8), rw, rh, [
        (title, 22, GREEN, True),
        ("", 6, GRAY, False),
        (price, 24, WHITE, True),
        ("", 6, GRAY, False),
        (desc, 15, GRAY, False),
    ])

add_text(s, "10 000 подписчиков = 9.9 млн ₸/мес",
         Inches(0), Inches(5.4), W, Inches(0.7),
         font_size=32, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s, "Только на подписках, без лидогенерации",
         Inches(0), Inches(6.1), W, Inches(0.5),
         font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(s, 8)

# ============================================================
# SLIDE 9 — TRACTION
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, "Что уже построено",
         Inches(0), Inches(0.5), W, Inches(0.8),
         font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Left: checklist
checks = [
    "Работающий MVP (фронтенд + бэкенд)",
    "Авторизация по номеру телефона",
    "Симулятор с реальными расчётами",
    "AI объяснение через Claude API",
    "Интерфейс на русском и казахском",
    "Тёмная и светлая тема",
]

txBox = s.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.5), Inches(5))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(checks):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
        p.space_before = Pt(14)
    p.text = f"✓  {item}"
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY
    p.font.name = "Calibri"
    # Make checkmark green via run
    p.clear()
    run1 = p.add_run()
    run1.text = "✓  "
    run1.font.size = Pt(20)
    run1.font.color.rgb = GREEN
    run1.font.bold = True
    run1.font.name = "Calibri"
    run2 = p.add_run()
    run2.text = item
    run2.font.size = Pt(18)
    run2.font.color.rgb = GRAY
    run2.font.name = "Calibri"

# Right: tech stack
add_text(s, "Tech Stack",
         Inches(7.5), Inches(1.8), Inches(4), Inches(0.5),
         font_size=16, color=GRAY, bold=True)

techs = ["React", "FastAPI", "PostgreSQL", "Claude AI", "Supabase"]
badge_x = Inches(7.5)
badge_y = Inches(2.6)
bw = Inches(2.2)
bh = Inches(0.65)
bgap = Inches(0.2)

for i, tech in enumerate(techs):
    row = i // 2
    col = i % 2
    x = badge_x + col * (bw + bgap)
    y = badge_y + row * (bh + bgap)
    card_with_text(s, x, y, bw, bh, [
        (tech, 15, GRAY, True),
    ], border=RGBColor(0x33, 0x33, 0x33))

add_slide_number(s, 9)

# ============================================================
# SLIDE 10 — CLOSING
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

closing = [
    ("10.2 млн", "казахстанцев берут кредиты"),
    ("0", "приложений показывающих реальный скор бесплатно"),
    ("1", "Nesie"),
]

cy_start = Inches(1.0)
row_height = Inches(2.0)

for i, (big, label) in enumerate(closing):
    y = cy_start + i * row_height
    add_text(s, big,
             Inches(0), y, W, Inches(1.2),
             font_size=72, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, label,
             Inches(0), y + Inches(1.0), W, Inches(0.5),
             font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

add_text(s, "Байқауға қатысушы: [ИМЯ]  |  Атырау, 2026",
         Inches(0), H - Inches(0.8), W, Inches(0.5),
         font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(s, 10)

# ============================================================
# SAVE
# ============================================================
out_path = os.path.join(os.path.dirname(__file__), "nesie_pitch.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
